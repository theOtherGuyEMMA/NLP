import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Search both notebooks' figures and pipeline results figures
FIG_DIRS = [Path('figures'), Path('results') / 'figures']
OUTPUT = 'Speech_Classification_Swahili.pptx'

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.level = 0
        p.text = b
        p.font.size = Pt(18)

def find_image(image_name: str) -> Path | None:
    for d in FIG_DIRS:
        p = d / image_name
        if p.exists():
            return p
    return None

def add_image_slide(prs, title, image_path, caption=None):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # Support passing either a Path or str name; if name only, search dirs
    img_path = Path(image_path) if isinstance(image_path, (str, Path)) else None
    if img_path and not img_path.exists():
        img_path = find_image(img_path.name)

    if img_path and img_path.exists():
        # Consistent sizing and centering
        max_width = Inches(9)
        left = Inches((10 - 9) / 2)  # center within 10" slide width
        top = Inches(1.5)
        pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
        if caption:
            txBox = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(9.5), Inches(1.2))
            tf = txBox.text_frame
            tf.text = caption
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.size = Pt(14)
    else:
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        tf = body.text_frame
        tf.text = f"[Figure not found: {img_path}]"

def build_presentation():
    prs = Presentation()

    # Title
    add_title_slide(
        prs,
        title="Swahili Speech Classification",
        subtitle="Preprocessing • Modeling (CNN/Wav2Vec2) • Evaluation • Robustness • Exploration"
    )

    # Problem & Objective
    add_bullets_slide(prs, "Problem & Objective", [
        "Task: Classify Swahili speech (keywords/intents).",
        "Datasets: Common Voice (Swahili) or Swahili Words Parallel Dataset.",
        "Deliverables: Notebook, Report, and Slide Deck with comparisons and plots."
    ])

    # Dataset Overview
    add_bullets_slide(prs, "Dataset Overview", [
        "Audio sampled near 16 kHz; variable durations.",
        "Labels: derive keyword classes or use folder-based labels.",
        "Inspect sample rate, duration, and label distribution."
    ])
    add_image_slide(prs, "Metadata Overview", 'metadata_overview.png', caption="Duration and label distributions")

    # Preprocessing
    add_bullets_slide(prs, "Preprocessing", [
        "Resample to 16 kHz; normalize features.",
        "Extract MFCC and Mel-spectrograms.",
        "Augment with noise and pitch shifts; optional time-stretch."
    ])
    add_image_slide(prs, "Waveform & Mel-spectrogram", 'waveform_mel_sample.png', caption="Example waveform with corresponding Mel-spectrogram")

    # Modeling
    add_bullets_slide(prs, "Modeling", [
        "1D CNN over MFCC/Mel features (Conv-ReLU-Pool).",
        "Wav2Vec2 fine-tuning (classification head).",
        "Experiment with conv layers, kernel sizes, activations."
    ])

    # Training & Regularization
    add_bullets_slide(prs, "Training & Regularization", [
        "Train 10–15 epochs; monitor train/val loss.",
        "Apply dropout, early stopping, and LR scheduling.",
        "Hyperparameter search on architecture and features."
    ])
    add_image_slide(prs, "CNN Training Losses", 'cnn1d_losses.png', caption="Training vs validation loss across epochs")

    # Evaluation
    add_bullets_slide(prs, "Evaluation", [
        "Metrics: Accuracy, Macro-F1, Confusion Matrix.",
        "Embedding visualization via PCA and t-SNE.",
        "Compare CNN vs Wav2Vec2 performance."
    ])
    add_image_slide(prs, "Confusion Matrix (CNN)", 'cnn1d_confusion.png', caption="Confusion matrix for CNN baseline")
    add_image_slide(prs, "Embeddings (PCA & t-SNE)", 'embeddings_pca_tsne.png', caption="2D projections of penultimate-layer embeddings")

    # Exploration
    add_bullets_slide(prs, "Exploration & Comparisons", [
        "Sampling rates: 8 kHz vs 16 kHz.",
        "Mel resolution: 40, 64, 128 bins.",
        "Attempt transfer learning from English Common Voice."
    ])
    add_image_slide(prs, "SR/Mel Comparison", 'sampling_mel_experiment.png', caption="Accuracy vs Mel bins across sampling rates")

    # Insert actual results from runner
    add_image_slide(prs, "Model Comparison", 'model_comparison.png', caption="Accuracy and training time across models")
    add_image_slide(prs, "Confusion Matrix (Random Forest)", 'confusion_matrix_random_forest.png', caption="Confusion matrix for Random Forest")
    add_image_slide(prs, "Confusion Matrix (SVM)", 'confusion_matrix_svm.png', caption="Confusion matrix for SVM")

    # Robustness & Generalization
    add_bullets_slide(prs, "Robustness & Generalization", [
        "Augmentations improve noise/pitch resilience.",
        "Speaker-wise splits for fair evaluation.",
        "Test-time augmentation and domain shifts (devices/environments)."
    ])

    # Extension pipeline
    add_bullets_slide(prs, "Extension: ASR + Sentiment Pipeline", [
        "ASR to transcribe Swahili (Whisper / Wav2Vec2).",
        "Optional translation (Sw → En).",
        "Sentiment model on text output."
    ])

    # Attribution & Licensing
    add_bullets_slide(prs, "Attribution & Licensing", [
        "Figures generated from experiments in this project.",
        "Dataset: Mozilla Common Voice (Swahili) — Clips licensed CC0 1.0.",
        "Ensure any external images used are properly licensed and attributed."
    ])

    # Conclusion
    add_bullets_slide(prs, "Conclusions & Next Steps", [
        "CNN is a solid baseline; Wav2Vec2 excels with GPU/data.",
        "Refine label schema and expand datasets.",
        "Deploy robust evaluation and production pipeline."
    ])

    prs.save(OUTPUT)
    return OUTPUT

if __name__ == '__main__':
    out = build_presentation()
    print(f"Saved: {out}")